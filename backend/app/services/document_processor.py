"""
Document Processing Pipeline
─────────────────────────────
Supports: PDF, DOCX, PPTX, XLSX, Images (with Tesseract OCR)

Pipeline per document:
  1. Detect file type
  2. Extract raw text (OCR for images/scanned PDFs)
  3. Clean & split text into overlapping chunks
  4. Generate embeddings (GPU if available)
  5. Store in FAISS index with metadata
"""
from __future__ import annotations
import asyncio
import logging
import re
import uuid
from pathlib import Path
from typing import List

import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import openpyxl

from app.core.config import settings
from app.core.vector_db import VectorDBManager
from app.models.document import DocumentRecord, DocumentStatus

logger = logging.getLogger(__name__)

pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD


# ── Text Extraction ───────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    """Extract text from PDF; fall back to OCR for scanned pages."""
    doc = fitz.open(str(path))
    pages: list[str] = []
    for page in doc:
        text = page.get_text("text").strip()
        if len(text) < 30:
            # Likely scanned — render as image and OCR
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img, config="--oem 3 --psm 6")
        pages.append(text)
    doc.close()
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    document = docx.Document(str(path))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"[Slide {i + 1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def _extract_image(path: Path) -> str:
    img = Image.open(str(path))
    return pytesseract.image_to_string(img, config="--oem 3 --psm 6")


EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xlsx,
    ".png":  _extract_image,
    ".jpg":  _extract_image,
    ".jpeg": _extract_image,
    ".tiff": _extract_image,
    ".bmp":  _extract_image,
}


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    extractor = EXTRACTORS.get(suffix)
    if not extractor:
        raise ValueError(f"Unsupported file type: {suffix}")
    text = extractor(path)
    # Basic cleaning
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


# ── Chunking ──────────────────────────────────────────────────────────────────

def split_into_chunks(text: str, chunk_size: int = settings.CHUNK_SIZE,
                      overlap: int = settings.CHUNK_OVERLAP) -> list[str]:
    """
    Split text on sentence/paragraph boundaries with a sliding window.
    Falls back to character split for very long strings.
    """
    # Split into sentences (simple heuristic)
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for sentence in sentences:
        s_len = len(sentence)
        if current_len + s_len > chunk_size and current:
            chunks.append(" ".join(current))
            # Overlap: keep the tail
            tail_len = 0
            tail: list[str] = []
            for s in reversed(current):
                if tail_len + len(s) > overlap:
                    break
                tail.insert(0, s)
                tail_len += len(s)
            current = tail
            current_len = tail_len

        current.append(sentence)
        current_len += s_len

    if current:
        chunks.append(" ".join(current))

    return [c for c in chunks if len(c.strip()) > 20]


# ── Main Pipeline ─────────────────────────────────────────────────────────────

async def process_document(record: DocumentRecord, file_path: Path) -> DocumentRecord:
    """
    Full pipeline: extract → chunk → embed → store.
    Updates `record` in-place and returns it.
    """
    record.status = DocumentStatus.PROCESSING
    try:
        logger.info("Processing %s (%s) …", record.original_name, record.file_type)

        # Run CPU-bound extraction in thread pool
        loop = asyncio.get_event_loop()
        raw_text = await loop.run_in_executor(None, extract_text, file_path)

        if not raw_text:
            raise ValueError("No text could be extracted from the document.")

        chunks = split_into_chunks(raw_text)
        logger.info("  → %d chunks generated.", len(chunks))

        metadata = [
            {
                "doc_id":    record.id,
                "filename":  record.original_name,
                "chunk_idx": i,
                "text":      chunk,
            }
            for i, chunk in enumerate(chunks)
        ]

        added = await VectorDBManager.add_chunks(chunks, metadata)
        record.chunk_count = added
        record.status = DocumentStatus.READY
        logger.info("  → %d vectors stored. Document ready.", added)

    except Exception as exc:
        logger.exception("Failed to process document %s: %s", record.id, exc)
        record.status = DocumentStatus.FAILED
        record.error = str(exc)

    return record
